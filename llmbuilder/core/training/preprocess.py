"""
Data preprocessing script for LLM training.
Extracts and cleans text from PDF, DOCX, and TXT files.
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional
from loguru import logger

# Document processing imports
try:
    import fitz  # PyMuPDF
except ImportError:
    logger.warning("PyMuPDF not available. PDF processing will be disabled.")
    fitz = None

try:
    import docx2txt
except ImportError:
    logger.warning("docx2txt not available. DOCX processing will be disabled.")
    docx2txt = None

try:
    import markdown
    from bs4 import BeautifulSoup
except ImportError:
    logger.warning("markdown/beautifulsoup4 not available. Enhanced markdown processing will be disabled.")
    markdown = None
    BeautifulSoup = None

try:
    from pptx import Presentation
except ImportError:
    logger.warning("python-pptx not available. PPTX processing will be disabled.")
    Presentation = None

# Removed sys.path.append - using proper package imports
from llmbuilder.core.training.utils import setup_logging
from llmbuilder.core.data.ingest import DocumentIngester
from llmbuilder.core.data.dedup import DeduplicationPipeline


class TextProcessor:
    """Handles text extraction and cleaning from various file formats."""
    
    def __init__(self, 
                 min_length: int = 100,
                 max_length: int = 1000000,
                 remove_urls: bool = True,
                 remove_emails: bool = True,
                 normalize_whitespace: bool = True):
        """
        Initialize text processor.
        
        Args:
            min_length: Minimum text length to keep
            max_length: Maximum text length to keep
            remove_urls: Whether to remove URLs
            remove_emails: Whether to remove email addresses
            normalize_whitespace: Whether to normalize whitespace
        """
        self.min_length = min_length
        self.max_length = max_length
        self.remove_urls = remove_urls
        self.remove_emails = remove_emails
        self.normalize_whitespace = normalize_whitespace
        
        # Compile regex patterns
        self.url_pattern = re.compile(
            r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\(\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
        )
        self.email_pattern = re.compile(
            r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        )
        self.whitespace_pattern = re.compile(r'\s+')
        
        # Additional cleaning patterns
        self.html_tag_pattern = re.compile(r'<[^>]+>')
        self.special_chars_pattern = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]')
        self.repeated_chars_pattern = re.compile(r'(.)\1{4,}')  # Remove 5+ repeated chars
        self.markdown_link_pattern = re.compile(r'\[([^\]]+)\]\([^\)]+\)')  # [text](url)
        self.markdown_image_pattern = re.compile(r'!\[([^\]]*)\]\([^\)]+\)')  # ![alt](url)
        
        logger.info("Text processor initialized")
    
    def extract_from_txt(self, file_path: Path) -> Optional[str]:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            return text
        except Exception as e:
            logger.error(f"Error reading TXT file {file_path}: {e}")
            return None
    
    def extract_from_pdf(self, file_path: Path) -> Optional[str]:
        """Extract text from PDF file with enhanced error handling."""
        if fitz is None:
            logger.warning(f"Skipping PDF {file_path}: PyMuPDF not available")
            return None
        
        try:
            # Try to open PDF with different methods
            doc = None
            text = ""
            
            try:
                doc = fitz.open(file_path)
            except Exception as e:
                logger.warning(f"Failed to open PDF {file_path} normally: {e}")
                # Try to open as stream for corrupted PDFs
                try:
                    with open(file_path, 'rb') as f:
                        doc = fitz.open(stream=f.read(), filetype="pdf")
                    logger.info(f"Successfully opened PDF {file_path} as stream")
                except Exception as e2:
                    logger.error(f"Failed to open PDF {file_path} as stream: {e2}")
                    return None
            
            if doc is None:
                return None
                
            # Extract text from all pages
            for page_num in range(doc.page_count):
                try:
                    page = doc[page_num]
                    page_text = page.get_text()
                    
                    # If regular text extraction fails, try OCR-like extraction
                    if not page_text.strip():
                        try:
                            # Try different text extraction methods
                            page_text = page.get_text("text", flags=11)  # More aggressive extraction
                        except:
                            page_text = ""
                    
                    text += page_text + "\n"
                    
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num} of {file_path}: {e}")
                    continue
            
            doc.close()
            
            if not text.strip():
                logger.warning(f"No text extracted from PDF {file_path}")
                return None
                
            return text
            
        except Exception as e:
            logger.error(f"Error reading PDF file {file_path}: {e}")
            return None
    
    def extract_from_docx(self, file_path: Path) -> Optional[str]:
        """Extract text from DOCX file."""
        if docx2txt is None:
            logger.warning(f"Skipping DOCX {file_path}: docx2txt not available")
            return None
        
        try:
            text = docx2txt.process(str(file_path))
            return text
        except Exception as e:
            logger.error(f"Error reading DOCX file {file_path}: {e}")
            return None
    
    def extract_from_md(self, file_path: Path) -> Optional[str]:
        """Extract text from Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # If markdown library is available, convert to HTML then extract text
            if markdown is not None and BeautifulSoup is not None:
                try:
                    # Convert markdown to HTML
                    html = markdown.markdown(content, extensions=['extra', 'codehilite'])
                    # Extract text from HTML
                    soup = BeautifulSoup(html, 'html.parser')
                    text = soup.get_text(separator=' ', strip=True)
                    return text
                except Exception as e:
                    logger.warning(f"Failed to process markdown with library for {file_path}: {e}")
                    # Fall back to basic processing
            
            # Basic markdown processing - remove common markdown syntax
            text = content
            
            # Remove markdown headers
            text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
            
            # Remove markdown links but keep text: [text](url) -> text
            text = self.markdown_link_pattern.sub(r'\1', text)
            
            # Remove markdown images
            text = self.markdown_image_pattern.sub('', text)
            
            # Remove markdown code blocks
            text = re.sub(r'```[\s\S]*?```', '', text)
            text = re.sub(r'`([^`]+)`', r'\1', text)  # Inline code
            
            # Remove markdown emphasis
            text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # Bold
            text = re.sub(r'\*([^*]+)\*', r'\1', text)  # Italic
            text = re.sub(r'__([^_]+)__', r'\1', text)  # Bold
            text = re.sub(r'_([^_]+)_', r'\1', text)  # Italic
            
            # Remove markdown lists
            text = re.sub(r'^[\s]*[-*+]\s+', '', text, flags=re.MULTILINE)
            text = re.sub(r'^[\s]*\d+\.\s+', '', text, flags=re.MULTILINE)
            
            # Remove horizontal rules
            text = re.sub(r'^[\s]*[-*_]{3,}[\s]*$', '', text, flags=re.MULTILINE)
            
            return text
            
        except Exception as e:
            logger.error(f"Error reading Markdown file {file_path}: {e}")
            return None
            
    def extract_from_pptx(self, file_path: Path) -> Optional[str]:
        """Extract text from PowerPoint (.pptx) file."""
        if Presentation is None:
            logger.warning(f"Skipping PPTX {file_path}: python-pptx not available")
            return None
            
        try:
            prs = Presentation(file_path)
            text_runs = []
            
            for slide in prs.slides:
                # Extract text from slide title
                if slide.shapes.title and slide.shapes.title.text:
                    text_runs.append(slide.shapes.title.text)
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                text_runs.append(run.text)
            
            # Join all text runs with newlines
            return '\n'.join(text_runs)
            
        except Exception as e:
            logger.error(f"Error reading PPTX file {file_path}: {e}")
            return None
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize text with enhanced cleaning."""
        if not text:
            return ""
        
        # Remove special control characters
        text = self.special_chars_pattern.sub('', text)
        
        # Remove HTML tags (in case any remain)
        text = self.html_tag_pattern.sub('', text)
        
        # Remove URLs
        if self.remove_urls:
            text = self.url_pattern.sub('', text)
        
        # Remove email addresses
        if self.remove_emails:
            text = self.email_pattern.sub('', text)
        
        # Remove repeated characters (like ===== or -----)
        text = self.repeated_chars_pattern.sub(r'\1\1\1', text)  # Keep max 3 repetitions
        
        # Normalize whitespace
        if self.normalize_whitespace:
            text = self.whitespace_pattern.sub(' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Remove lines that are mostly special characters or numbers
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if len(line) < 3:  # Skip very short lines
                continue
            # Skip lines that are mostly numbers, punctuation, or special chars
            if len(re.sub(r'[\d\s\W]', '', line)) < len(line) * 0.3:
                continue
            cleaned_lines.append(line)
        
        text = '\n'.join(cleaned_lines)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def process_file(self, file_path: Path) -> Optional[str]:
        """Process a single file and extract clean text."""
        logger.info(f"Processing {file_path}")
        
        # Extract text based on file extension
        text = None
        ext = file_path.suffix.lower()
        
        if ext == '.txt':
            text = self.extract_from_txt(file_path)
        elif ext == '.pdf':
            text = self.extract_from_pdf(file_path)
        elif ext == '.docx':
            text = self.extract_from_docx(file_path)
        elif ext == '.md':
            text = self.extract_from_md(file_path)
        elif ext == '.pptx':
            text = self.extract_from_pptx(file_path)
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return None
        
        if text is None:
            logger.warning(f"Failed to extract text from {file_path}")
            return None
        
        # Clean the text
        cleaned_text = self.clean_text(text)
        
        # Check length constraints
        if len(cleaned_text) < self.min_length:
            logger.warning(f"Text too short ({len(cleaned_text)} chars): {file_path}")
            return None
        
        if len(cleaned_text) > self.max_length:
            logger.warning(f"Text too long ({len(cleaned_text)} chars), truncating: {file_path}")
            cleaned_text = cleaned_text[:self.max_length]
        
        logger.info(f"Extracted {len(cleaned_text)} characters from {file_path}")
        return cleaned_text


class DataPreprocessor:
    """Main data preprocessing pipeline with integrated ingestion and deduplication."""
    
    def __init__(self, 
                 raw_data_dir: str = "data/raw",
                 cleaned_data_dir: str = "data/cleaned",
                 deduped_data_dir: str = "data/deduped",
                 use_enhanced_ingestion: bool = True,
                 use_deduplication: bool = True,
                 config: Optional[Dict[str, Any]] = None,
                 **processor_kwargs):
        """
        Initialize data preprocessor.
        
        Args:
            raw_data_dir: Directory containing raw data files
            cleaned_data_dir: Directory to save cleaned data
            deduped_data_dir: Directory to save deduplicated data
            use_enhanced_ingestion: Whether to use enhanced ingestion pipeline
            use_deduplication: Whether to use deduplication pipeline
            config: Configuration dictionary
            **processor_kwargs: Arguments for TextProcessor
        """
        self.raw_data_dir = Path(raw_data_dir)
        self.cleaned_data_dir = Path(cleaned_data_dir)
        self.deduped_data_dir = Path(deduped_data_dir)
        self.use_enhanced_ingestion = use_enhanced_ingestion
        self.use_deduplication = use_deduplication
        self.config = config or {}
        
        # Initialize components
        self.processor = TextProcessor(**processor_kwargs)
        
        # Initialize enhanced ingestion if enabled
        if self.use_enhanced_ingestion:
            try:
                ingestion_config = self.config.get('ingestion', {})
                self.ingester = DocumentIngester(
                    output_dir=str(self.cleaned_data_dir),
                    max_file_size_mb=ingestion_config.get('max_file_size_mb', 100)
                )
                logger.info("Enhanced ingestion pipeline initialized")
            except Exception as e:
                logger.warning(f"Failed to initialize enhanced ingestion: {e}")
                self.use_enhanced_ingestion = False
                self.ingester = None
        else:
            self.ingester = None
        
        # Initialize deduplication if enabled
        if self.use_deduplication:
            try:
                dedup_config = self.config.get('deduplication', {})
                embedding_config = dedup_config.get('embedding_deduplication', {})
                self.deduplicator = DeduplicationPipeline(
                    use_hash_dedup=dedup_config.get('hash_deduplication', {}).get('enabled', True),
                    use_embedding_dedup=embedding_config.get('enabled', True),
                    embedding_model=embedding_config.get('model_name', 'all-MiniLM-L6-v2'),
                    similarity_threshold=embedding_config.get('similarity_threshold', 0.85)
                )
                logger.info("Deduplication pipeline initialized")
            except Exception as e:
                logger.warning(f"Failed to initialize deduplication: {e}")
                self.use_deduplication = False
                self.deduplicator = None
        else:
            self.deduplicator = None
        
        # Create output directories
        self.cleaned_data_dir.mkdir(parents=True, exist_ok=True)
        if self.use_deduplication:
            self.deduped_data_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"Data preprocessor initialized")
        logger.info(f"Raw data dir: {self.raw_data_dir}")
        logger.info(f"Cleaned data dir: {self.cleaned_data_dir}")
        if self.use_deduplication:
            logger.info(f"Deduped data dir: {self.deduped_data_dir}")
        logger.info(f"Enhanced ingestion: {self.use_enhanced_ingestion}")
        logger.info(f"Deduplication: {self.use_deduplication}")
    
    def find_files(self) -> List[Path]:
        """Find all supported files in raw data directory."""
        if not self.raw_data_dir.exists():
            raise FileNotFoundError(f"Raw data directory not found: {self.raw_data_dir}")
        
        supported_extensions = ['.txt', '.pdf', '.docx', '.md', '.pptx']
        files = []
        
        for ext in supported_extensions:
            files.extend(self.raw_data_dir.glob(f"**/*{ext}"))
        
        logger.info(f"Found {len(files)} files to process")
        return sorted(files)
    
    def get_unprocessed_files(self, files: List[Path]) -> List[Path]:
        """Filter out files that have already been processed."""
        unprocessed = []
        skipped = 0
        
        for file_path in files:
            output_file = self.cleaned_data_dir / f"{file_path.stem}_cleaned.txt"
            
            # Check if output file exists and is newer than source
            if output_file.exists():
                try:
                    source_mtime = file_path.stat().st_mtime
                    output_mtime = output_file.stat().st_mtime
                    
                    if output_mtime >= source_mtime:
                        # File already processed and up to date
                        skipped += 1
                        logger.debug(f"Skipping {file_path.name} (already processed)")
                        continue
                    else:
                        logger.debug(f"Processing {file_path.name} (source newer than output)")
                except OSError as e:
                    # If we can't check timestamps, process the file
                    logger.debug(f"Cannot check timestamps for {file_path.name}: {e}")
                    pass
            else:
                logger.debug(f"Processing {file_path.name} (no output file)")
            
            unprocessed.append(file_path)
        
        logger.info(f"File filtering: {skipped} already processed, {len(unprocessed)} need processing")
        
        # Debug: Show first few files and their status
        if len(files) > 0:
            sample_file = files[0]
            output_file = self.cleaned_data_dir / f"{sample_file.stem}_cleaned.txt"
            logger.info(f"Debug sample - Source: {sample_file}")
            logger.info(f"Debug sample - Output: {output_file}")
            logger.info(f"Debug sample - Output exists: {output_file.exists()}")
            if output_file.exists():
                source_mtime = sample_file.stat().st_mtime
                output_mtime = output_file.stat().st_mtime
                logger.info(f"Debug sample - Source mtime: {source_mtime}")
                logger.info(f"Debug sample - Output mtime: {output_mtime}")
                logger.info(f"Debug sample - Output newer: {output_mtime >= source_mtime}")
        return unprocessed
    
    def process_all(self) -> Dict[str, Any]:
        """Process all files in the raw data directory with integrated pipeline."""
        from tqdm import tqdm
        import traceback
        
        logger.info("Starting integrated data preprocessing pipeline...")
        logger.info(f"Looking for files in: {self.raw_data_dir.absolute()}")
        
        try:
            # Step 1: Enhanced Ingestion (if enabled)
            if self.use_enhanced_ingestion and self.ingester:
                logger.info("=== Step 1: Enhanced Document Ingestion ===")
                try:
                    ingestion_results = self.ingester.ingest_directory(
                        self.raw_data_dir, 
                        recursive=True
                    )
                    logger.info(f"Ingestion complete: {ingestion_results['processed_count']} files processed")
                    
                    # Use ingested files for further processing
                    processed_files = list(self.cleaned_data_dir.glob("*.txt"))
                    
                except Exception as e:
                    logger.error(f"Enhanced ingestion failed: {e}")
                    logger.info("Falling back to legacy processing")
                    self.use_enhanced_ingestion = False
                    processed_files = []
            else:
                processed_files = []
            
            # Step 2: Legacy Processing (if enhanced ingestion disabled or failed)
            if not self.use_enhanced_ingestion or not processed_files:
                logger.info("=== Step 2: Legacy Document Processing ===")
                all_files = self.find_files()
                logger.info(f"Found {len(all_files)} total files")
                
                # Filter out already processed files
                files = self.get_unprocessed_files(all_files)
                already_processed = len(all_files) - len(files)
                
                if already_processed > 0:
                    logger.info(f"Skipping {already_processed} already processed files")
                
                if not files and already_processed == 0:
                    logger.warning("No files found to process")
                    return {
                        'processed': 0,
                        'failed': 0,
                        'total_chars': 0,
                        'output_files': []
                    }
                
                if files:
                    logger.info(f"Processing {len(files)} new/modified files")
                    processed_files = self._process_files_legacy(files)
                else:
                    processed_files = list(self.cleaned_data_dir.glob("*_cleaned.txt"))
            
            # Step 3: Deduplication (if enabled)
            final_files = processed_files
            if self.use_deduplication and self.deduplicator and processed_files:
                logger.info("=== Step 3: Deduplication ===")
                try:
                    # Convert Path objects to strings for deduplication
                    file_paths = [str(f) for f in processed_files if f.suffix == '.txt']
                    
                    if file_paths:
                        dedup_stats = self.deduplicator.process_files(
                            file_paths, 
                            str(self.deduped_data_dir)
                        )
                        logger.info(f"Deduplication complete: {dedup_stats.final_document_count} files remaining")
                        logger.info(f"Removed {dedup_stats.exact_duplicates_removed} exact duplicates")
                        logger.info(f"Removed {dedup_stats.near_duplicates_removed} near duplicates")
                        
                        # Use deduplicated files as final output
                        final_files = list(self.deduped_data_dir.glob("*.txt"))
                    else:
                        logger.warning("No text files found for deduplication")
                        
                except Exception as e:
                    logger.error(f"Deduplication failed: {e}")
                    logger.info("Using non-deduplicated files")
            
            # Step 4: Create combined file and calculate statistics
            logger.info("=== Step 4: Finalizing Output ===")
            return self._finalize_output(final_files)
        
            processed_count = 0
            failed_count = 0
            total_chars = 0
            output_files = []
            
            # Store the count of already processed files for final reporting
            already_processed_count = already_processed
            
            # Initialize tqdm progress bar
            progress_bar = tqdm(
                files, 
                desc="Processing files", 
                unit="file",
                bar_format='{l_bar}{bar:50}{r_bar}{bar:-10b}',
                colour='green'
            )
            
            for file_path in progress_bar:
                try:
                    # Process file
                    progress_bar.set_postfix(file=file_path.name[:20] + '...' if len(file_path.name) > 20 else file_path.name)
                    logger.debug(f"Processing file: {file_path}")
                    
                    cleaned_text = self.processor.process_file(file_path)
                    
                    if cleaned_text:
                        # Ensure output directory exists
                        self.cleaned_data_dir.mkdir(parents=True, exist_ok=True)
                        
                        # Save cleaned text to file
                        output_file = self.cleaned_data_dir / f"{file_path.stem}_cleaned.txt"
                        try:
                            with open(output_file, 'w', encoding='utf-8') as f:
                                f.write(cleaned_text)
                            
                            processed_count += 1
                            total_chars += len(cleaned_text)
                            output_files.append(str(output_file.absolute()))
                            progress_bar.write(f"[DONE] Processed: {file_path.name} -> {len(cleaned_text):,} chars")
                        except Exception as write_error:
                            failed_count += 1
                            progress_bar.write(f"[FAIL] Failed to write {file_path.name}: {str(write_error)}")
                            logger.error(f"Error writing {output_file}: {traceback.format_exc()}")
                    else:
                        failed_count += 1
                        progress_bar.write(f"[SKIP] No content extracted from: {file_path.name}")
                
                except Exception as e:
                    failed_count += 1
                    error_msg = f"✗ Error processing {file_path.name}: {str(e)}"
                    # Remove any non-ASCII characters from error message
                    safe_error_msg = error_msg.encode('ascii', 'replace').decode('ascii')
                    progress_bar.write(safe_error_msg)
                    logger.error(f"Error processing {file_path}: {traceback.format_exc()}")
        
            # Always ensure combined file exists with all cleaned files
            combined_path = self.cleaned_data_dir / "combined_text.txt"
            all_cleaned_files = list(self.cleaned_data_dir.glob("*_cleaned.txt"))
            
            if all_cleaned_files:
                try:
                    logger.info(f"Creating/updating combined file from {len(all_cleaned_files)} cleaned files")
                    with open(combined_path, 'w', encoding='utf-8') as combined_file:
                        for cleaned_file in sorted(all_cleaned_files):
                            try:
                                with open(cleaned_file, 'r', encoding='utf-8') as f:
                                    content = f.read().strip()
                                    if content:  # Only write non-empty content
                                        combined_file.write(content)
                                        combined_file.write("\n\n")  # Separator between documents
                            except Exception as e:
                                logger.warning(f"Failed to read {cleaned_file} for combining: {e}")
                                continue
                    
                    # Verify the combined file has content
                    if combined_path.exists() and combined_path.stat().st_size > 0:
                        logger.info(f"Combined file ready: {combined_path} ({combined_path.stat().st_size:,} bytes)")
                    else:
                        logger.warning("Combined file is empty")
                        
                except Exception as e:
                    logger.error(f"Failed to create combined file: {e}")
                    logger.error(traceback.format_exc())
            
            # Include already processed files in the final count
            all_output_files = list(self.cleaned_data_dir.glob("*_cleaned.txt"))
            if self.cleaned_data_dir.joinpath("combined_text.txt").exists():
                all_output_files.append(self.cleaned_data_dir / "combined_text.txt")
            
            # Calculate total characters from all files
            total_all_chars = 0
            for output_file in all_output_files:
                try:
                    with open(output_file, 'r', encoding='utf-8') as f:
                        total_all_chars += len(f.read())
                except:
                    pass
            
            results = {
                'processed': processed_count + already_processed_count,
                'failed': failed_count,
                'total_chars': total_all_chars,
                'output_files': [str(f.absolute()) for f in all_output_files]
            }
            
            if processed_count > 0:
                logger.info(f"Preprocessing complete: {processed_count} newly processed, {already_processed_count} already up-to-date, {failed_count} failed")
            else:
                logger.info(f"All {already_processed_count} files were already up-to-date")
            logger.info(f"Total characters: {total_all_chars:,}")
            
            # Only raise error if we had files to process but none succeeded and some failed
            if processed_count == 0 and failed_count > 0 and files and already_processed_count == 0:
                error_msg = "No files were successfully processed. Check the logs for details."
                logger.error(error_msg)
                raise RuntimeError(error_msg)
                
            return results
            
        except Exception as e:
            logger.error(f"Fatal error in preprocessing: {e}")
            logger.error(traceback.format_exc())
            raise
    
    def _process_files_legacy(self, files: List[Path]) -> List[Path]:
        """Process files using legacy method."""
        from tqdm import tqdm
        import traceback
        
        processed_files = []
        failed_count = 0
        
        # Initialize tqdm progress bar
        progress_bar = tqdm(
            files, 
            desc="Processing files", 
            unit="file",
            bar_format='{l_bar}{bar:50}{r_bar}{bar:-10b}',
            colour='green'
        )
        
        for file_path in progress_bar:
            try:
                # Process file
                progress_bar.set_postfix(file=file_path.name[:20] + '...' if len(file_path.name) > 20 else file_path.name)
                logger.debug(f"Processing file: {file_path}")
                
                cleaned_text = self.processor.process_file(file_path)
                
                if cleaned_text:
                    # Ensure output directory exists
                    self.cleaned_data_dir.mkdir(parents=True, exist_ok=True)
                    
                    # Save cleaned text to file
                    output_file = self.cleaned_data_dir / f"{file_path.stem}_cleaned.txt"
                    try:
                        with open(output_file, 'w', encoding='utf-8') as f:
                            f.write(cleaned_text)
                        
                        processed_files.append(output_file)
                        progress_bar.write(f"[DONE] Processed: {file_path.name} -> {len(cleaned_text):,} chars")
                    except Exception as write_error:
                        failed_count += 1
                        progress_bar.write(f"[FAIL] Failed to write {file_path.name}: {str(write_error)}")
                        logger.error(f"Error writing {output_file}: {traceback.format_exc()}")
                else:
                    failed_count += 1
                    progress_bar.write(f"[SKIP] No content extracted from: {file_path.name}")
            
            except Exception as e:
                failed_count += 1
                error_msg = f"✗ Error processing {file_path.name}: {str(e)}"
                # Remove any non-ASCII characters from error message
                safe_error_msg = error_msg.encode('ascii', 'replace').decode('ascii')
                progress_bar.write(safe_error_msg)
                logger.error(f"Error processing {file_path}: {traceback.format_exc()}")
        
        logger.info(f"Legacy processing complete: {len(processed_files)} processed, {failed_count} failed")
        return processed_files
    
    def _finalize_output(self, final_files: List[Path]) -> Dict[str, Any]:
        """Finalize output and create combined file."""
        import traceback
        
        # Determine the final output directory
        if self.use_deduplication and final_files and any(self.deduped_data_dir.name in str(f) for f in final_files):
            output_dir = self.deduped_data_dir
            combined_name = "combined_deduped.txt"
        else:
            output_dir = self.cleaned_data_dir
            combined_name = "combined_text.txt"
        
        # Create combined file
        combined_path = output_dir / combined_name
        total_chars = 0
        
        if final_files:
            try:
                logger.info(f"Creating combined file from {len(final_files)} files")
                with open(combined_path, 'w', encoding='utf-8') as combined_file:
                    for file_path in sorted(final_files):
                        try:
                            with open(file_path, 'r', encoding='utf-8') as f:
                                content = f.read().strip()
                                if content:  # Only write non-empty content
                                    combined_file.write(content)
                                    combined_file.write("\n\n")  # Separator between documents
                                    total_chars += len(content)
                        except Exception as e:
                            logger.warning(f"Failed to read {file_path} for combining: {e}")
                            continue
                
                # Verify the combined file has content
                if combined_path.exists() and combined_path.stat().st_size > 0:
                    logger.info(f"Combined file ready: {combined_path} ({combined_path.stat().st_size:,} bytes)")
                else:
                    logger.warning("Combined file is empty")
                    
            except Exception as e:
                logger.error(f"Failed to create combined file: {e}")
                logger.error(traceback.format_exc())
        
        # Prepare output file list
        output_files = [str(f.absolute()) for f in final_files]
        if combined_path.exists():
            output_files.append(str(combined_path.absolute()))
        
        # Calculate total characters if not already done
        if total_chars == 0:
            for file_path in final_files:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        total_chars += len(f.read())
                except:
                    pass
        
        results = {
            'processed': len(final_files),
            'failed': 0,  # Failed count handled in individual steps
            'total_chars': total_chars,
            'output_files': output_files,
            'pipeline_used': {
                'enhanced_ingestion': self.use_enhanced_ingestion,
                'deduplication': self.use_deduplication,
                'final_output_dir': str(output_dir)
            }
        }
        
        logger.info(f"Pipeline complete: {len(final_files)} files, {total_chars:,} total characters")
        logger.info(f"Final output directory: {output_dir}")
        
        return results


def load_config(config_path: str = "config.json") -> Dict[str, Any]:
    """Load configuration from file."""
    import json
    
    try:
        with open(config_path, 'r') as f:
            config = json.load(f)
        logger.info(f"Configuration loaded from {config_path}")
        return config
    except FileNotFoundError:
        logger.warning(f"Configuration file {config_path} not found, using defaults")
        return {}
    except Exception as e:
        logger.error(f"Error loading configuration: {e}")
        return {}


def main():
    """Main entry point with integrated pipeline support."""
    import argparse
    
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="LLMBuilder Data Preprocessing Pipeline")
    parser.add_argument('--config', type=str, default='config.json', help='Configuration file path')
    parser.add_argument('--no-ingestion', action='store_true', help='Disable enhanced ingestion')
    parser.add_argument('--no-dedup', action='store_true', help='Disable deduplication')
    parser.add_argument('--raw-dir', type=str, help='Raw data directory')
    parser.add_argument('--cleaned-dir', type=str, help='Cleaned data directory')
    parser.add_argument('--deduped-dir', type=str, help='Deduplicated data directory')
    
    args = parser.parse_args()
    
    # Setup logging
    setup_logging(log_dir="logs", level="INFO")
    
    try:
        # Load configuration
        config = load_config(args.config)
        
        # Get directories from config or args
        paths = config.get('paths', {})
        raw_dir = args.raw_dir or paths.get('raw_data_dir', 'data/raw')
        cleaned_dir = args.cleaned_dir or paths.get('cleaned_data_dir', 'data/cleaned')
        deduped_dir = args.deduped_dir or paths.get('deduped_data_dir', 'data/deduped')
        
        # Get preprocessing settings
        preprocessing = config.get('preprocessing', {})
        
        # Initialize preprocessor with integrated pipeline
        preprocessor = DataPreprocessor(
            raw_data_dir=raw_dir,
            cleaned_data_dir=cleaned_dir,
            deduped_data_dir=deduped_dir,
            use_enhanced_ingestion=not args.no_ingestion,
            use_deduplication=not args.no_dedup,
            config=config,
            min_length=preprocessing.get('min_length', 50),
            max_length=preprocessing.get('max_length', 500000),
            remove_urls=True,
            remove_emails=True,
            normalize_whitespace=preprocessing.get('normalize_whitespace', True)
        )
        
        # Process all files
        results = preprocessor.process_all()
        
        # Print summary
        logger.info("=== Integrated Pipeline Summary ===")
        logger.info(f"Files processed: {results['processed']}")
        logger.info(f"Files failed: {results['failed']}")
        logger.info(f"Total characters: {results['total_chars']:,}")
        logger.info(f"Output files: {len(results['output_files'])}")
        
        # Print pipeline information
        pipeline_info = results.get('pipeline_used', {})
        logger.info(f"Enhanced ingestion: {pipeline_info.get('enhanced_ingestion', False)}")
        logger.info(f"Deduplication: {pipeline_info.get('deduplication', False)}")
        logger.info(f"Final output: {pipeline_info.get('final_output_dir', 'N/A')}")
        
        if results['processed'] == 0:
            logger.warning("No files were successfully processed!")
            logger.info("Please check that you have files in the data/raw directory")
            logger.info("Supported formats: .txt, .pdf, .docx, .md, .html, .epub")
        
    except Exception as e:
        logger.error(f"Preprocessing failed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()

